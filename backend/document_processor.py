"""
Enhanced Multi-format document processing module
Supports: PDF, Word, PowerPoint, Excel, Audio, Video, Images, Text, HTML, Markdown, RTF, CSV
"""
import os
import pandas as pd
import pdfplumber
import pytesseract
from PIL import Image
import whisper
import speech_recognition as sr
from pathlib import Path
from typing import Dict, List, Optional, Union
import tempfile
import subprocess
import json

# Document processing imports
from docx import Document
from pptx import Presentation
import mammoth
from bs4 import BeautifulSoup
import markdown
from striprtf.striprtf import rtf_to_text
import xml.etree.ElementTree as ET
import zipfile
import ebooklib
from ebooklib import epub


class DocumentProcessor:
    def __init__(self):
        self.whisper_model = None
        self.recognizer = sr.Recognizer()
        self.supported_formats = {
            'pdf': ['.pdf'],
            'word': ['.docx', '.doc'],
            'powerpoint': ['.pptx', '.ppt'],
            'excel': ['.xlsx', '.xls', '.csv'],
            'audio': ['.mp3', '.wav', '.m4a', '.flac', '.ogg', '.aac'],
            'video': ['.mp4', '.avi', '.mov', '.mkv', '.flv', '.wmv'],
            'image': ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp'],
            'text': ['.txt', '.log'],
            'html': ['.html', '.htm'],
            'markdown': ['.md', '.markdown'],
            'rtf': ['.rtf'],
            'json': ['.json'],
            'xml': ['.xml'],
            'epub': ['.epub'],
            'code': ['.py', '.js', '.java', '.cpp', '.c', '.cs', '.php', '.rb', '.go', '.rs', '.swift']
        }
    
    def detect_file_type(self, file_path: str) -> str:
        """
        Automatically detect file type from extension
        """
        # Handle both file paths and uploaded file objects
        if hasattr(file_path, 'name'):
            # Streamlit UploadedFile object
            ext = Path(file_path.name).suffix.lower()
        else:
            # String file path
            ext = Path(file_path).suffix.lower()
        
        for file_type, extensions in self.supported_formats.items():
            if ext in extensions:
                return file_type
        
        raise ValueError(f"Unsupported file extension: {ext}")
    
    def process_document(self, file_path: str, file_type: Optional[str] = None, password: str = None) -> Dict:
        """
        Main entry point for document processing
        Auto-detects file type if not provided
        
        Args:
            file_path: Path to the file
            file_type: Optional file type override
            password: Optional password for encrypted PDFs
        """
        if file_type is None:
            file_type = self.detect_file_type(file_path)
        
        processors = {
            'pdf': lambda fp: self.process_pdf(fp, password=password),  # Pass password here
            'word': self.process_word,
            'powerpoint': self.process_powerpoint,
            'excel': self.process_excel,
            'audio': self.process_audio,
            'video': self.process_video,
            'image': self.process_image,
            'text': self.process_text,
            'html': self.process_html,
            'markdown': self.process_markdown,
            'rtf': self.process_rtf,
            'json': self.process_json,
            'xml': self.process_xml,
            'epub': self.process_epub,
            'code': self.process_code
        }
        
        if file_type not in processors:
            raise ValueError(f"Unsupported file type: {file_type}")
        
        return processors[file_type](file_path)
    
    def process_pdf(self, file_path: str, password: str = None) -> Dict:
        """
        Process PDF files with text extraction and OCR fallback
        Handles password-protected PDFs
        """

        text_content = []
        metadata = {
            'file_type': 'pdf',
            'file_name': Path(file_path).name,
            'pages': 0,
            'encrypted': False
        }
        
        try:
            # Attempt to open PDF with password
            try:
                with pdfplumber.open(file_path, password=password or "") as pdf:
                    metadata['pages'] = len(pdf.pages)
                    
                    for page_num, page in enumerate(pdf.pages, 1):
                        text = page.extract_text()
                        
                        if text and text.strip():
                            text_content.append({
                                'page': page_num,
                                'content': text.strip(),
                                'type': 'text'
                            })
                        else:
                            try:
                                img = page.to_image(resolution=300)
                                pil_img = img.original
                                ocr_text = pytesseract.image_to_string(pil_img)
                                
                                if ocr_text.strip():
                                    text_content.append({
                                        'page': page_num,
                                        'content': ocr_text.strip(),
                                        'type': 'ocr'
                                    })
                            except Exception as ocr_error:
                                print(f"OCR failed for page {page_num}: {ocr_error}")
            
            except Exception as pdf_error:
                # Check if it's specifically a password error
                error_type = type(pdf_error).__name__
                error_str = str(pdf_error)
                
                if "PDFPasswordIncorrect" in error_type or "password" in error_str.lower():
                    metadata['encrypted'] = True
                    # Raise a user-friendly error
                    raise ValueError(
                        "🔒 PASSWORD-PROTECTED PDF: This PDF file is encrypted and requires a password. "
                        "Please unlock the PDF before uploading or provide the correct password."
                    )
                else:
                    # Re-raise other PDF errors
                    raise
        
        except ValueError as ve:
            # Re-raise ValueError (our password error) as-is
            raise Exception(str(ve))
        
        except Exception as e:
            raise Exception(f"Error processing PDF: {str(e)}")
        
        return {
            'metadata': metadata,
            'content': text_content,
            'full_text': ' '.join([item['content'] for item in text_content])
        }
        
    def process_word(self, file_path: str) -> Dict:
        """
        Process Word documents (.docx, .doc)
        """
        text_content = []
        metadata = {
            'file_type': 'word',
            'file_name': Path(file_path).name
        }
        
        try:
            file_ext = Path(file_path).suffix.lower()
            
            if file_ext == '.docx':
                doc = Document(file_path)
                
                # Extract paragraphs
                for i, para in enumerate(doc.paragraphs):
                    if para.text.strip():
                        text_content.append({
                            'paragraph': i + 1,
                            'content': para.text.strip(),
                            'style': para.style.name
                        })
                
                # Extract tables
                for i, table in enumerate(doc.tables):
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    
                    text_content.append({
                        'table': i + 1,
                        'content': table_data,
                        'type': 'table'
                    })
                
                metadata['paragraphs'] = len(doc.paragraphs)
                metadata['tables'] = len(doc.tables)
            
            elif file_ext == '.doc':
                # Use mammoth for .doc files
                with open(file_path, "rb") as docx_file:
                    result = mammoth.extract_raw_text(docx_file)
                    text_content.append({
                        'content': result.value,
                        'type': 'text'
                    })
        
        except Exception as e:
            raise Exception(f"Error processing Word document: {str(e)}")
        
        full_text_parts = []
        for item in text_content:
            if 'content' in item:
                if isinstance(item['content'], str):
                    full_text_parts.append(item['content'])
                elif isinstance(item['content'], list):
                    for row in item['content']:
                        full_text_parts.append(' | '.join(row))
        
        return {
            'metadata': metadata,
            'content': text_content,
            'full_text': '\n'.join(full_text_parts)
        }
    
    def process_powerpoint(self, file_path: str) -> Dict:
        """
        Process PowerPoint presentations (.pptx, .ppt)
        """
        text_content = []
        metadata = {
            'file_type': 'powerpoint',
            'file_name': Path(file_path).name
        }
        
        try:
            prs = Presentation(file_path)
            metadata['slides'] = len(prs.slides)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    
                    # Extract text from tables
                    if hasattr(shape, "table"):
                        table_data = []
                        for row in shape.table.rows:
                            row_data = [cell.text.strip() for cell in row.cells]
                            table_data.append(row_data)
                        slide_text.append(str(table_data))
                
                if slide_text:
                    text_content.append({
                        'slide': slide_num,
                        'content': '\n'.join(slide_text),
                        'layout': slide.slide_layout.name
                    })
        
        except Exception as e:
            raise Exception(f"Error processing PowerPoint: {str(e)}")
        
        return {
            'metadata': metadata,
            'content': text_content,
            'full_text': '\n\n'.join([item['content'] for item in text_content])
        }
    
    def process_excel(self, file_path: str) -> Dict:
        """
        Process Excel files and CSV
        """
        file_ext = Path(file_path).suffix.lower()
        
        try:
            if file_ext == '.csv':
                df = pd.read_csv(file_path)
            else:
                df = pd.read_excel(file_path, sheet_name=None)
                
                if isinstance(df, dict):
                    all_data = []
                    for sheet_name, sheet_df in df.items():
                        sheet_data = {
                            'sheet_name': sheet_name,
                            'data': sheet_df.to_dict('records'),
                            'columns': [str(col) for col in sheet_df.columns],
                            'summary': sheet_df.describe().to_dict() if not sheet_df.empty else {}
                        }
                        all_data.append(sheet_data)
                    
                    text_content = []
                    for sheet in all_data:
                        text_content.append(f"Sheet: {sheet['sheet_name']}")
                        text_content.append(f"Columns: {', '.join(sheet['columns'])}")
                        text_content.append(str(pd.DataFrame(sheet['data'])))
                    
                    return {
                        'metadata': {
                            'file_type': 'excel',
                            'file_name': Path(file_path).name,
                            'sheets': len(all_data)
                        },
                        'content': all_data,
                        'full_text': '\n\n'.join(text_content)
                    }
            
            return {
                'metadata': {
                    'file_type': 'excel',
                    'file_name': Path(file_path).name,
                    'rows': len(df),
                    'columns': [str(col) for col in df.columns]
                },
                'content': df.to_dict('records'),
                'full_text': df.to_string()
            }
        
        except Exception as e:
            raise Exception(f"Error processing Excel: {str(e)}")
    
    def process_audio(self, file_path: str) -> Dict:
        """
        Transcribe audio using Whisper
        """
        try:
            if self.whisper_model is None:
                self.whisper_model = whisper.load_model("base")
            
            result = self.whisper_model.transcribe(file_path)
            
            return {
                'metadata': {
                    'file_type': 'audio',
                    'file_name': Path(file_path).name,
                    'language': result.get('language', 'unknown')
                },
                'content': result.get('segments', []),
                'full_text': result['text']
            }
        
        except Exception as e:
            raise Exception(f"Error processing audio: {str(e)}")
    
    def process_video(self, file_path: str) -> Dict:
        """
        Extract audio from video and transcribe
        """
        try:
            with tempfile.NamedTemporaryFile(suffix='.wav', delete=False) as tmp_audio:
                audio_path = tmp_audio.name
            
            cmd = [
                'ffmpeg', '-i', file_path,
                '-vn', '-acodec', 'pcm_s16le',
                '-ar', '16000', '-ac', '1',
                audio_path, '-y'
            ]
            
            try:
                subprocess.run(cmd, check=True, capture_output=True)
            except subprocess.CalledProcessError:
                raise Exception("ffmpeg not found. Please install ffmpeg for video processing.")
            
            audio_result = self.process_audio(audio_path)
            os.unlink(audio_path)
            
            audio_result['metadata']['file_type'] = 'video'
            audio_result['metadata']['file_name'] = Path(file_path).name
            
            return audio_result
        
        except Exception as e:
            raise Exception(f"Error processing video: {str(e)}")
    
    def process_image(self, file_path: str) -> Dict:
        """
        OCR text extraction from images
        """
        try:
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            
            return {
                'metadata': {
                    'file_type': 'image',
                    'file_name': Path(file_path).name,
                    'dimensions': image.size,
                    'format': image.format
                },
                'content': [{
                    'type': 'ocr',
                    'content': text.strip()
                }],
                'full_text': text.strip()
            }
        except Exception as e:
            raise Exception(f"Error processing image: {str(e)}")
    
    def process_text(self, file_path: str) -> Dict:
        """
        Process plain text files
        """
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
            
            return {
                'metadata': {
                    'file_type': 'text',
                    'file_name': Path(file_path).name,
                    'size_bytes': os.path.getsize(file_path)
                },
                'content': [{'content': text}],
                'full_text': text
            }
        except Exception as e:
            raise Exception(f"Error processing text file: {str(e)}")
    
    def process_html(self, file_path: str) -> Dict:
        """
        Process HTML files
        """
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                html_content = f.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            text = soup.get_text()
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)
            
            return {
                'metadata': {
                    'file_type': 'html',
                    'file_name': Path(file_path).name,
                    'title': soup.title.string if soup.title else None
                },
                'content': [{
                    'html': html_content,
                    'text': text
                }],
                'full_text': text
            }
        except Exception as e:
            raise Exception(f"Error processing HTML: {str(e)}")
    
    def process_markdown(self, file_path: str) -> Dict:
        """
        Process Markdown files
        """
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                md_content = f.read()
            
            html = markdown.markdown(md_content)
            soup = BeautifulSoup(html, 'html.parser')
            text = soup.get_text()
            
            return {
                'metadata': {
                    'file_type': 'markdown',
                    'file_name': Path(file_path).name
                },
                'content': [{
                    'markdown': md_content,
                    'html': html,
                    'text': text
                }],
                'full_text': text
            }
        except Exception as e:
            raise Exception(f"Error processing Markdown: {str(e)}")
    
    def process_rtf(self, file_path: str) -> Dict:
        """
        Process RTF files
        """
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                rtf_content = f.read()
            
            text = rtf_to_text(rtf_content)
            
            return {
                'metadata': {
                    'file_type': 'rtf',
                    'file_name': Path(file_path).name
                },
                'content': [{'content': text}],
                'full_text': text
            }
        except Exception as e:
            raise Exception(f"Error processing RTF: {str(e)}")
    
    def process_json(self, file_path: str) -> Dict:
        """
        Process JSON files
        """
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                json_data = json.load(f)
            
            text = json.dumps(json_data, indent=2)
            
            return {
                'metadata': {
                    'file_type': 'json',
                    'file_name': Path(file_path).name
                },
                'content': json_data,
                'full_text': text
            }
        except Exception as e:
            raise Exception(f"Error processing JSON: {str(e)}")
    
    def process_xml(self, file_path: str) -> Dict:
        """
        Process XML files
        """
        try:
            tree = ET.parse(file_path)
            root = tree.getroot()
            
            def xml_to_dict(element):
                result = {}
                if element.attrib:
                    result['@attributes'] = element.attrib
                if element.text and element.text.strip():
                    result['text'] = element.text.strip()
                for child in element:
                    child_data = xml_to_dict(child)
                    if child.tag in result:
                        if not isinstance(result[child.tag], list):
                            result[child.tag] = [result[child.tag]]
                        result[child.tag].append(child_data)
                    else:
                        result[child.tag] = child_data
                return result
            
            data = xml_to_dict(root)
            text = ET.tostring(root, encoding='unicode', method='text')
            
            return {
                'metadata': {
                    'file_type': 'xml',
                    'file_name': Path(file_path).name,
                    'root_tag': root.tag
                },
                'content': data,
                'full_text': text
            }
        except Exception as e:
            raise Exception(f"Error processing XML: {str(e)}")
    
    def process_epub(self, file_path: str) -> Dict:
        """
        Process EPUB files
        """
        try:
            book = epub.read_epub(file_path)
            text_content = []
            
            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), 'html.parser')
                    text = soup.get_text()
                    if text.strip():
                        text_content.append(text.strip())
            
            return {
                'metadata': {
                    'file_type': 'epub',
                    'file_name': Path(file_path).name,
                    'title': book.get_metadata('DC', 'title'),
                    'author': book.get_metadata('DC', 'creator')
                },
                'content': text_content,
                'full_text': '\n\n'.join(text_content)
            }
        except Exception as e:
            raise Exception(f"Error processing EPUB: {str(e)}")
    
    def process_code(self, file_path: str) -> Dict:
        """
        Process source code files
        """
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                code = f.read()
            
            ext = Path(file_path).suffix.lower()
            language_map = {
                '.py': 'Python', '.js': 'JavaScript', '.java': 'Java',
                '.cpp': 'C++', '.c': 'C', '.cs': 'C#', '.php': 'PHP',
                '.rb': 'Ruby', '.go': 'Go', '.rs': 'Rust', '.swift': 'Swift'
            }
            
            return {
                'metadata': {
                    'file_type': 'code',
                    'file_name': Path(file_path).name,
                    'language': language_map.get(ext, 'Unknown'),
                    'lines': len(code.splitlines())
                },
                'content': [{'code': code}],
                'full_text': code
            }
        except Exception as e:
            raise Exception(f"Error processing code file: {str(e)}")
    
    def get_supported_formats(self) -> Dict[str, List[str]]:
        """
        Return all supported file formats
        """
        return self.supported_formats
    
    def extract_text_from_image(self, image_path: str) -> str:
        """
        OCR text extraction from images (legacy method for compatibility)
        """
        try:
            image = Image.open(image_path)
            text = pytesseract.image_to_string(image)
            return text.strip()
        except Exception as e:
            raise Exception(f"Error extracting text from image: {str(e)}")


# Example usage
if __name__ == "__main__":
    processor = DocumentProcessor()
    
    # Display supported formats
    print("Supported formats:")
    for file_type, extensions in processor.get_supported_formats().items():
        print(f"{file_type}: {', '.join(extensions)}")
    
    # Process a document
    # result = processor.process_document('example.pdf')
    # print(result['full_text'])